from pptx import Presentation
from pptx.util import Inches
from typing import List

class PPTService:
    @staticmethod
    def create_presentation(image_paths: List[str], output_file: str):
        """
        Creates a PPTX file where each slide is a full-screen image.
        """
        prs = Presentation()
        # Set to 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6] 
        
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save(output_file)
        print(f"✅ PowerPoint saved to: {output_file}")